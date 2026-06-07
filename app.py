import os
import re
import hmac
import subprocess
import tempfile
import urllib.request
from pathlib import Path
from datetime import date
from collections import OrderedDict

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# 画面最上部に表示する更新情報
VERSION = "## Version-2 — 20260607 000000 更新 - 簡易ログイン機能を追加"

# 設定変数

# スライドサイズ
SLIDE_W, SLIDE_H = Cm(33.87), Cm(19.05)

# テーブル列幅 [種目, Set, Rep, 重量]
COL_W = [Cm(5.66), Cm(1.31), Cm(1.31), Cm(2.61)]
TABLE_W = sum(COL_W)

# 配置
COLS = 3
MAX_PER_PAGE = 9
MARGIN_L = Cm(0.35)
COL_GAP = (SLIDE_W - MARGIN_L * 2 - TABLE_W * COLS) // (COLS - 1)
ROW_Y = [Cm(1.65), Cm(7.13), Cm(12.76)]

# 各パーツの高さ
TOP_BAND_H = Cm(0.12)
NAME_H = Cm(0.55)
HDR_H = Cm(0.37)
DATA_H = Cm(0.46)

# セル内余白
CELL_V = Cm(0.085)
CELL_H = Cm(0.10)

# 色（# なし 6桁 hex）
C_RED = "DC143C"
C_HDR_BG = "EEEEEE"
C_HDR_FG = "333333"
C_ODD = "FFFFFF"
C_EVEN = "F2F2F2"
C_BORDER = "CCCCCC"
C_DATE = "BBBBBB"
C_WHITE = "FFFFFF"

# フォントサイズ (pt)
FS_TITLE = 16
FS_NAME = 10
FS_TABLE = 6.5
FS_DATE = 9

# フォント。Linux(GCP)には Meiryo が無いため環境変数で差し替え可能にする
# ローカルWindows: 既定の Meiryo / GCP(Cloud Run): Noto Sans CJK JP を環境変数で指定
FONT = os.environ.get("PPTX_FONT", "Meiryo")
RECT = 1  # msoShapeRectangle


# 列名・列範囲の計算

def num_to_col_name(n: int) -> str:
    """1ベースの列番号をスプレッドシートのアルファベット列名に変換（1->A, 27->AA）"""
    name = ""
    while n > 0:
        n, remainder = divmod(n - 1, 26)
        name = chr(65 + remainder) + name
    return name


def get_date_col_range(target_date_str: str) -> str:
    """入力日付(MM/DD)から取得する列範囲(例 'L:O')を返す"""
    base_date_str = "04/22"
    base_start_col = 12  # 'L'
    cols_per_day = 4
    target_date = pd.Timestamp(f"2024/{target_date_str}")
    base_date = pd.Timestamp(f"2024/{base_date_str}")
    delta_days = (target_date - base_date).days
    start_col_num = base_start_col + (delta_days * cols_per_day)
    end_col_num = start_col_num + cols_per_day - 1
    return f"{num_to_col_name(start_col_num)}:{num_to_col_name(end_col_num)}"


# データ取得（常に最新スプレッドシート）

def env_value(key: str, default: str = "") -> str:
    """設定値を 環境変数 -> .env の順で取得"""
    v = os.environ.get(key)
    if v:
        return v
    env_path = Path(__file__).parent / ".env"
    if env_path.exists():
        for line in env_path.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            k, val = line.split("=", 1)
            if k.strip() == key:
                return val.strip().strip('"').strip("'")
    return default


def get_sheet_id() -> str:
    """GOOGLE_SHEET_ID を取得"""
    return env_value("GOOGLE_SHEET_ID")


def download_spreadsheet(dest_dir: str) -> Path:
    """Googleスプレッドシートを xlsx でダウンロードして保存先パスを返す"""
    sid = get_sheet_id()
    if not sid:
        raise RuntimeError("GOOGLE_SHEET_ID が未設定です（.env またはデプロイ環境変数に設定してください）")
    url = f"https://docs.google.com/spreadsheets/d/{sid}/export?format=xlsx"
    dest = Path(dest_dir) / "workout.xlsx"
    urllib.request.urlretrieve(url, dest)
    # xlsx の先頭4バイトは PK\x03\x04（ZIP）。HTMLが返る＝非公開等の失敗
    with open(dest, "rb") as f:
        if f.read(4) != b"PK\x03\x04":
            dest.unlink(missing_ok=True)
            raise RuntimeError(
                "ダウンロード失敗：シートが非公開か SHEET_ID が誤っています。"
                "共有設定を『リンクを知っている全員が閲覧可』にしてください"
            )
    return dest


def extract_specific_date(xlsx_path: Path, date_extract: str) -> pd.DataFrame:
    """該当日の全選手データを1つのDataFrameに集約（CSVを介さず直接返す）"""
    target_cols = get_date_col_range(date_extract)
    xls = pd.ExcelFile(xlsx_path)
    frames = []
    for sheet in xls.sheet_names:
        try:
            df = pd.read_excel(xls, sheet_name=sheet, usecols=target_cols,
                               skiprows=2, nrows=99, header=None)
        except ValueError:
            # 指定列範囲がシートの列数を超える場合はスキップ
            continue
        if df.empty:
            continue
        df.columns = ["種目", "セット", "レップ", "重量"]
        df.dropna(subset=["セット", "レップ", "重量"], how="all", inplace=True)
        if df.empty:
            continue
        df.insert(0, "選手名", sheet)
        frames.append(df)
    if not frames:
        raise RuntimeError(
            f"日付 '{date_extract}'（計算列範囲 '{target_cols}'）に該当データがありません"
        )
    return pd.concat(frames).fillna("")


def df_to_players(df: pd.DataFrame) -> OrderedDict:
    """DataFrame を {選手名: [(種目, Set, Rep, 重量), ...]} に変換"""
    players = OrderedDict()
    for _, r in df.iterrows():
        name = str(r["選手名"]).strip()
        if not name:
            continue
        players.setdefault(name, []).append(
            (str(r["種目"]).strip(), str(r["セット"]).strip(),
             str(r["レップ"]).strip(), str(r["重量"]).strip())
        )
    return players


# PPTX 描画

def rgb(h: str) -> RGBColor:
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def fmt_num(val: str) -> str:
    """数値なら整数表示（3.0->3）、文字列はそのまま返す"""
    try:
        f = float(val)
        return str(int(f)) if f == int(f) else val
    except (ValueError, TypeError):
        return val


def solid_rect(slide, x, y, w, h, color: str):
    """塗りつぶし矩形（罫線なし）"""
    s = slide.shapes.add_shape(RECT, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(color)
    s.line.fill.background()


def add_text(slide, x, y, w, h, text, fs, bold=False, color=C_DATE,
             align=PP_ALIGN.LEFT, mt=0, ml=Cm(0.05), bg=None):
    """テキストボックスを追加。bg 指定で背景色付き"""
    tx = slide.shapes.add_textbox(x, y, w, h)
    if bg:
        tx.fill.solid()
        tx.fill.fore_color.rgb = rgb(bg)
    tx.line.fill.background()
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_top, tf.margin_bottom = mt, 0
    tf.margin_left, tf.margin_right = ml, 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = FONT


def _apply_cell(cell, bg: str, text, fs, bold=False, fg=C_HDR_FG, align=PP_ALIGN.LEFT):
    """セル背景・罫線・テキストを一括設定"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr")).set("val", bg)
    w_val = str(int(Pt(0.4).emu / 12700))
    for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = etree.SubElement(tcPr, qn(side))
        ln.set("w", w_val)
        sf2 = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf2, qn("a:srgbClr")).set("val", C_BORDER)
    tf = cell.text_frame
    tf.margin_top = tf.margin_bottom = CELL_V
    tf.margin_left = tf.margin_right = CELL_H
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text) if text else ""
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = rgb(fg)
    run.font.name = FONT

    # 空白セルのデフォルトフォント（Calibri 18pt）を上書き
    para_el = p._p
    for old in para_el.findall(qn("a:endParaRPr")):
        para_el.remove(old)
    epr = etree.SubElement(para_el, qn("a:endParaRPr"))
    epr.set("sz", str(int(fs * 100)))
    if bold:
        epr.set("b", "1")
    etree.SubElement(epr, qn("a:latin")).set("typeface", FONT)


def draw_player(slide, name, exercises, x, y):
    """選手名行＋ヘッダー行＋データ行を1つのテーブルで描画"""
    n_data = len(exercises)
    n_rows = 1 + n_data
    tbl_h = NAME_H + HDR_H + n_data * DATA_H

    table = slide.shapes.add_table(n_rows, 4, x, y, TABLE_W, tbl_h).table

    for ci, cw in enumerate(COL_W):
        table.columns[ci].width = cw

    table.rows[0].height = NAME_H
    for ri in range(1, n_rows):
        table.rows[ri].height = DATA_H

    # ヘッダー行（選手名 + 列見出し）
    for ci, hdr in enumerate([name, "Set", "Rep", "重量"]):
        _apply_cell(table.cell(0, ci), C_RED, hdr, FS_NAME, fg=C_WHITE, bold=True)

    # データ行
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT]
    for ri, (item, sets, reps, weight) in enumerate(exercises):
        bg = C_ODD if ri % 2 == 0 else C_EVEN
        row = [item, fmt_num(sets), fmt_num(reps), fmt_num(weight)]
        for ci, (val, al) in enumerate(zip(row, aligns)):
            _apply_cell(table.cell(ri + 1, ci), bg, val, FS_TABLE, fg="1A1A1A", align=al)

    return tbl_h


def build_slide(prs, page_players, title_text, date_text):
    """1スライドを生成"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("FFFFFF")

    solid_rect(slide, 0, 0, SLIDE_W, TOP_BAND_H, C_RED)
    add_text(slide, MARGIN_L, TOP_BAND_H, Cm(12), ROW_Y[0] - TOP_BAND_H,
             title_text, FS_TITLE, bold=True, color=C_RED, mt=Cm(0.25))
    add_text(slide, SLIDE_W - Cm(6), TOP_BAND_H, Cm(6) - MARGIN_L, ROW_Y[0] - TOP_BAND_H,
             date_text, FS_DATE, color=C_DATE, align=PP_ALIGN.RIGHT, mt=Cm(0.35))

    col_x = [MARGIN_L + i * (TABLE_W + COL_GAP) for i in range(COLS)]
    for idx, (name, exercises) in enumerate(page_players):
        draw_player(slide, name, exercises, col_x[idx % COLS], ROW_Y[idx // COLS])


def build_pptx(players_items, date_extract: str, out_dir: str) -> Path:
    """全選手データから pptx を生成して保存先パスを返す"""
    m = re.search(r"(\d{2})/(\d{2})", date_extract)
    month, day = (int(m.group(1)), int(m.group(2))) if m else (date.today().month, date.today().day)
    today = date.today()
    title_text = f"{month}月{day}日 トレーニング実施記録"
    date_text = f"{today.year}/{today.month:02d}/{today.day:02d} Workout Log"

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    for i in range(0, len(players_items), MAX_PER_PAGE):
        build_slide(prs, players_items[i:i + MAX_PER_PAGE], title_text, date_text)

    out = Path(out_dir) / f"{month}月{day}日トレーニング実施記録.pptx"
    prs.save(str(out))
    return out


def pptx_to_pdf(pptx_path: Path, out_dir: str) -> Path:
    """LibreOffice(soffice)で pptx -> pdf 変換"""
    subprocess.run(
        ["soffice", "--headless",
         "-env:UserInstallation=file:///tmp/lo_profile",
         "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
        check=True, timeout=180,
    )
    pdf = Path(out_dir) / (Path(pptx_path).stem + ".pdf")
    if not pdf.exists():
        raise RuntimeError("PDF変換に失敗しました（LibreOffice未インストールの可能性）")
    return pdf


# 簡易ログイン

def check_login() -> bool:
    """ログイン済みなら True。未ログインならログイン画面を表示して False を返す"""
    if st.session_state.get("logged_in"):
        return True

    st.caption(VERSION)
    st.title("ログイン")
    user_correct = env_value("APP_USERNAME")
    pw_correct = env_value("APP_PASSWORD")

    in_user = st.text_input("ユーザー名")
    in_pw = st.text_input("パスワード", type="password")

    if st.button("ログイン", type="primary", use_container_width=True):
        if not user_correct or not pw_correct:
            st.error("認証情報（APP_USERNAME / APP_PASSWORD）が未設定です")
        elif hmac.compare_digest(in_user, user_correct) and hmac.compare_digest(in_pw, pw_correct):
            st.session_state["logged_in"] = True
            st.rerun()
        else:
            st.error("ユーザー名またはパスワードが違います")
    return False


# Streamlit UI（スマホからボタン1つで実行）

def main():
    st.set_page_config(page_title="トレーニング記録PDF生成", page_icon="🏋️")

    if not check_login():
        return

    st.caption(VERSION)
    st.title("トレーニング実施記録 PDF 生成")
    st.write("最新のスプレッドシートを取得し、選択した実施日の記録を PDF で出力します。")

    if st.button("ログアウト"):
        st.session_state["logged_in"] = False
        st.rerun()

    target = st.date_input("実施日を選択", value=date.today())

    if st.button("PDFを生成", type="primary", use_container_width=True):
        date_extract = target.strftime("%m/%d")
        try:
            with st.spinner("生成中…（スプレッドシート取得 → 集計 → PDF変換）"):
                with tempfile.TemporaryDirectory() as tmp:
                    xlsx = download_spreadsheet(tmp)
                    df = extract_specific_date(xlsx, date_extract)
                    players = list(df_to_players(df).items())
                    pptx = build_pptx(players, date_extract, tmp)
                    pdf = pptx_to_pdf(pptx, tmp)
                    pdf_bytes = pdf.read_bytes()
                    pdf_name = pdf.name
            st.success(f"完了：{len(players)}名分を出力しました")
            st.download_button("PDFをダウンロード", data=pdf_bytes,
                               file_name=pdf_name, mime="application/pdf",
                               use_container_width=True)
        except Exception as e:
            st.error(str(e))


if __name__ == "__main__":
    main()
